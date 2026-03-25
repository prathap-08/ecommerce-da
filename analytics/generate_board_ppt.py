from pathlib import Path
import json
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt


REPORTS_DIR = Path("reports")
SCREEN_DIR = Path("dashboard/screenshots")
OUT_FILE = REPORTS_DIR / "board_summary.pptx"


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullets_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.font.size = Pt(18)


def add_image_slide(prs: Presentation, title: str, image_path: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    if image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(0.7), Inches(1.3), width=Inches(12.0), height=Inches(5.4))
    else:
        tx = slide.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(10), Inches(1))
        tx.text_frame.text = f"Image not found: {image_path.name}"


def main() -> None:
    kpis = json.loads((REPORTS_DIR / "kpis.json").read_text(encoding="utf-8"))
    monthly = pd.read_csv(REPORTS_DIR / "monthly_trend.csv", parse_dates=["order_date"])
    profit_cat = pd.read_csv(REPORTS_DIR / "profit_by_category.csv")
    forecast = pd.read_csv(REPORTS_DIR / "sales_forecast.csv", parse_dates=["month"])
    rfm_summary = pd.read_csv(REPORTS_DIR / "crm_rfm_summary.csv") if (REPORTS_DIR / "crm_rfm_summary.csv").exists() else pd.DataFrame()

    top_month = monthly.sort_values("revenue", ascending=False).iloc[0]
    low_month = monthly.sort_values("revenue", ascending=True).iloc[0]
    top_cat = profit_cat.sort_values("profit", ascending=False).iloc[0]
    low_margin_cat = profit_cat.sort_values("profit_margin_pct", ascending=True).iloc[0]
    peak_fc = forecast.sort_values("forecast_revenue", ascending=False).iloc[0]

    prs = Presentation()

    add_title_slide(
        prs,
        "Board Summary: E-Commerce Analytics",
        "Generated from integrated data, SQL, Python analytics, and ML outputs",
    )

    add_bullets_slide(
        prs,
        "Executive KPI Snapshot",
        [
            f"Revenue: {kpis['total_revenue']:,.2f}",
            f"Profit: {kpis['total_profit']:,.2f}",
            f"Margin: {kpis['profit_margin_pct']:.2f}%",
            f"Delivered Orders: {kpis['total_orders']:,}",
            f"Average Order Value: {kpis['avg_order_value']:,.2f}",
        ],
    )

    add_bullets_slide(
        prs,
        "Demand and Profit Insights",
        [
            f"Peak month: {top_month['order_date'].strftime('%Y-%m')} ({top_month['revenue']:,.2f})",
            f"Lowest month: {low_month['order_date'].strftime('%Y-%m')} ({low_month['revenue']:,.2f})",
            f"Top profit category: {top_cat['category']} ({top_cat['profit']:,.2f})",
            f"Lowest margin category: {low_margin_cat['category']} ({low_margin_cat['profit_margin_pct']:.2f}%)",
        ],
    )

    if not rfm_summary.empty:
        top_seg = rfm_summary.sort_values("revenue", ascending=False).iloc[0]
        add_bullets_slide(
            prs,
            "Customer Segmentation (RFM)",
            [
                f"Highest value segment: {top_seg['rfm_segment']}",
                f"Segment customers: {int(top_seg['customers']):,}",
                f"Segment revenue: {top_seg['revenue']:,.2f}",
                "Use CRM priority tags P1/P2/P3 from export to sequence campaigns.",
            ],
        )

    add_bullets_slide(
        prs,
        "Forecast and Actions",
        [
            f"Forecasted peak month: {peak_fc['month'].strftime('%Y-%m')} ({peak_fc['forecast_revenue']:,.2f})",
            "Scale inventory and logistics 6-8 weeks ahead of projected peaks.",
            "Activate cross-sell recommendations for high-frequency cohorts.",
            "Protect margins by reviewing deep-discount categories monthly.",
        ],
    )

    add_image_slide(prs, "Monthly Revenue Trend", SCREEN_DIR / "monthly_revenue_trend.png")
    add_image_slide(prs, "Category Profit", SCREEN_DIR / "category_profit.png")
    add_image_slide(prs, "Retention Curve", SCREEN_DIR / "retention_curve.png")

    OUT_FILE.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_FILE)
    print(f"Generated {OUT_FILE}")


if __name__ == "__main__":
    main()
